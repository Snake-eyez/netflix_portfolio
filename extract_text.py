import os
from pptx import Presentation
from pypdf import PdfReader

def extract_pptx(filepath):
    text_content = []
    try:
        prs = Presentation(filepath)
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            if slide_text:
                text_content.append(f"--- Slide {i+1} ---\n" + "\n".join(slide_text))
        return "\n".join(text_content)
    except Exception as e:
        return f"Error reading PPTX: {str(e)}"

def extract_pdf(filepath):
    text_content = []
    try:
        reader = PdfReader(filepath)
        for i, page in enumerate(reader.pages):
            text = page.extract_text()
            if text:
                text_content.append(f"--- Page {i+1} ---\n{text}")
        return "\n".join(text_content)
    except Exception as e:
        return f"Error reading PDF: {str(e)}"

folder = "presentation"
files = [
    "Presentatie van prototype.pptx",
    "ai_work_flowgrammers_multi_agent_20251106150313.pptx",
    "ai_workshop_hands_on.pdf"
]

with open("presentation_content.txt", "w", encoding="utf-8") as f:
    for filename in files:
        filepath = os.path.join(folder, filename)
        f.write(f"\n\n{'='*50}\nFILE: {filename}\n{'='*50}\n")
        if filename.endswith(".pptx"):
            f.write(extract_pptx(filepath))
        elif filename.endswith(".pdf"):
            f.write(extract_pdf(filepath))

print("Extraction complete.")
